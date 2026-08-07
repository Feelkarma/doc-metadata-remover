"""Tests for the metadata-removal core.

These tests build small Office documents with known metadata (using
python-docx / python-pptx / openpyxl, which are dev-only dependencies), run the
cleaner, and assert that (a) metadata is gone and (b) content is preserved.

Run with:  pytest -q
"""

import datetime
import os
import sys
import zipfile

import pytest

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from metadata_remover import (  # noqa: E402
    clean_document,
    legacy_support_available,
    pdf_support_available,
)

FIXTURES = os.path.join(os.path.dirname(__file__), "fixtures")
_OLE_META_STREAMS = ("\x05SummaryInformation", "\x05DocumentSummaryInformation")

docx = pytest.importorskip("docx")
pptx = pytest.importorskip("pptx")
openpyxl = pytest.importorskip("openpyxl")


def _read(path, part):
    with zipfile.ZipFile(path) as z:
        return z.read(part).decode("utf-8", "ignore") if part in z.namelist() else None


def test_docx_metadata_removed_and_content_preserved(tmp_path):
    from docx import Document

    src = tmp_path / "in.docx"
    d = Document()
    d.add_paragraph("Body text that must survive.")
    cp = d.core_properties
    cp.author = "Secret Author"
    cp.last_modified_by = "Secret Editor"
    cp.title = "Secret Title"
    cp.revision = 9
    cp.created = datetime.datetime(2020, 1, 1)
    d.save(str(src))

    result = clean_document(str(src), output_dir=str(tmp_path))
    assert result.success, result.error
    assert os.path.exists(result.output_path)
    assert os.path.exists(str(src))  # original preserved

    core = _read(result.output_path, "docProps/core.xml")
    for leaked in ("Secret Author", "Secret Editor", "Secret Title", "2020"):
        assert leaked not in core

    out = Document(result.output_path)
    assert [p.text for p in out.paragraphs] == ["Body text that must survive."]
    assert out.core_properties.author in ("", None)
    assert out.core_properties.last_modified_by in ("", None)


def _build_docx_with_comment(path):
    """Write a minimal .docx that has a classic comment and a tracked change."""
    parts = {
        "[Content_Types].xml":
            b'<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            b'<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            b'<Default Extension="xml" ContentType="application/xml"/>'
            b'<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            b'<Override PartName="/word/comments.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.comments+xml"/>'
            b'</Types>',
        "_rels/.rels":
            b'<?xml version="1.0"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            b'<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>',
        "word/_rels/document.xml.rels":
            b'<?xml version="1.0"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            b'<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments" Target="comments.xml"/></Relationships>',
        "word/document.xml":
            b'<?xml version="1.0"?><w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body>'
            b'<w:p><w:commentRangeStart w:id="0"/><w:r><w:t>Hello world</w:t></w:r><w:commentRangeEnd w:id="0"/>'
            b'<w:r><w:commentReference w:id="0"/></w:r></w:p>'
            b'<w:p><w:ins w:id="1" w:author="Jane Editor" w:date="2020-01-01T00:00:00Z"><w:r><w:t>inserted text</w:t></w:r></w:ins></w:p>'
            b'</w:body></w:document>',
        "word/comments.xml":
            b'<?xml version="1.0"?><w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            b'<w:comment w:id="0" w:author="Secret Reviewer" w:date="2020-01-01T00:00:00Z" w:initials="SR">'
            b'<w:p><w:r><w:t>This is my comment text</w:t></w:r></w:p></w:comment></w:comments>',
    }
    with zipfile.ZipFile(path, "w") as z:
        for name, blob in parts.items():
            z.writestr(name, blob)


def test_comments_and_tracked_changes_kept_but_anonymized(tmp_path):
    src = tmp_path / "in.docx"
    _build_docx_with_comment(str(src))

    result = clean_document(str(src), output_dir=str(tmp_path))
    assert result.success, result.error

    comments = _read(result.output_path, "word/comments.xml")
    document = _read(result.output_path, "word/document.xml")

    # The comment part must still exist and keep its text.
    assert comments is not None, "comment part was removed (should be kept)"
    assert "This is my comment text" in comments
    # ...but the reviewer identity must be gone.
    assert "Secret Reviewer" not in comments
    assert "SR" not in comments
    assert "2020" not in comments

    # Tracked change stays visible but the author is anonymised.
    assert "<w:ins" in document
    assert "inserted text" in document
    assert "Jane Editor" not in document
    # The comment anchor stays so the comment remains attached to the text.
    assert "commentReference" in document


def test_pptx_metadata_removed_and_content_preserved(tmp_path):
    from pptx import Presentation

    src = tmp_path / "in.pptx"
    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[0])
    slide.shapes.title.text = "Deck Title"
    p.core_properties.author = "Deck Author"
    p.save(str(src))

    result = clean_document(str(src), output_dir=str(tmp_path))
    assert result.success, result.error

    core = _read(result.output_path, "docProps/core.xml")
    assert "Deck Author" not in core

    out = Presentation(result.output_path)
    titles = [
        sh.text_frame.text
        for s in out.slides
        for sh in s.shapes
        if sh.has_text_frame
    ]
    assert "Deck Title" in titles


def test_xlsx_metadata_removed_and_content_preserved(tmp_path):
    from openpyxl import Workbook, load_workbook

    src = tmp_path / "in.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Header"
    ws["A2"] = 42
    ws["B2"] = "=A2*2"
    wb.properties.creator = "XL Creator"
    wb.properties.title = "XL Secret"
    wb.save(str(src))

    result = clean_document(str(src), output_dir=str(tmp_path))
    assert result.success, result.error

    core = _read(result.output_path, "docProps/core.xml")
    assert "XL Creator" not in core and "XL Secret" not in core

    out = load_workbook(result.output_path)
    ws2 = out.active
    assert ws2["A1"].value == "Header"
    assert ws2["A2"].value == 42
    assert ws2["B2"].value == "=A2*2"


def test_unsupported_extension(tmp_path):
    src = tmp_path / "note.txt"
    src.write_text("hello")
    result = clean_document(str(src), output_dir=str(tmp_path))
    assert not result.success
    assert "Unsupported" in (result.error or "")


@pytest.mark.skipif(not pdf_support_available(), reason="pikepdf not installed")
def test_pdf_metadata_removed_and_content_preserved(tmp_path):
    import pikepdf

    src = tmp_path / "in.pdf"
    pdf = pikepdf.new()
    pdf.add_blank_page(page_size=(200, 200))
    with pdf.open_metadata() as meta:
        meta["dc:creator"] = ["Secret PDF Author"]
        meta["dc:title"] = "Secret PDF Title"
    pdf.docinfo["/Author"] = "Secret PDF Author"
    pdf.docinfo["/Title"] = "Secret PDF Title"
    pdf.save(str(src))
    pdf.close()

    result = clean_document(str(src), output_dir=str(tmp_path))
    assert result.success, result.error
    assert os.path.exists(str(src))  # original preserved

    out = pikepdf.open(result.output_path)
    try:
        assert len(out.pages) == 1  # content preserved
        assert len(out.docinfo) == 0  # doc info stripped
        assert "/Metadata" not in out.Root  # XMP stream removed
    finally:
        out.close()


legacy = pytest.importorskip("olefile")


def _ole_streams(path):
    """Return {stream_name: bytes} for every stream in an OLE2 file."""
    import olefile

    out = {}
    ole = olefile.OleFileIO(path)
    try:
        for entry in ole.listdir():
            name = "/".join(entry)
            out[name] = ole.openstream(entry).read()
    finally:
        ole.close()
    return out


@pytest.mark.skipif(
    not legacy_support_available(), reason="olefile not installed"
)
@pytest.mark.parametrize("fixture", ["sample.doc", "sample.ppt", "sample.xls"])
def test_legacy_metadata_removed_and_content_preserved(tmp_path, fixture):
    import olefile

    src = os.path.join(FIXTURES, fixture)
    assert os.path.exists(src), f"missing fixture {fixture}"

    result = clean_document(src, output_dir=str(tmp_path))
    assert result.success, result.error
    assert result.output_path.lower().endswith(os.path.splitext(fixture)[1])
    assert os.path.exists(src)  # original preserved

    # Metadata must be gone from the standard property streams.
    ole = olefile.OleFileIO(result.output_path)
    try:
        meta = ole.get_metadata()
        assert meta.author in (None, b"", b"\x00")
        assert meta.title in (None, b"", b"\x00")
        assert meta.last_saved_by in (None, b"", b"\x00")
    finally:
        ole.close()

    # Every NON-metadata stream must be byte-for-byte identical (content and
    # formatting preserved); the two metadata streams keep their size but change.
    before = _ole_streams(src)
    after = _ole_streams(result.output_path)
    assert set(before) == set(after)
    for name in before:
        if name in _OLE_META_STREAMS:
            assert len(before[name]) == len(after[name])  # same size, blanked
        else:
            assert before[name] == after[name], f"content stream {name} changed"


def test_original_never_overwritten(tmp_path):
    from docx import Document

    src = tmp_path / "keep.docx"
    Document().save(str(src))
    result = clean_document(str(src), output_dir=str(tmp_path))
    assert result.success
    assert os.path.abspath(result.output_path) != os.path.abspath(str(src))
